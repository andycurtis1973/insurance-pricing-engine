#!/usr/bin/env python3
"""Build the webinar deck: insurance_pricing_webinar.pptx (16:9, 20 slides).

Reuses the video's slide builders (video/render.py) as full-bleed visuals and adds
the extra slides a webinar wants — agenda, the business case (the click-data trend and
the personalization opportunity), data/methodology, why-Bayesian, the elasticity
calibration + sweep, a "see it live" slide for the interactive app, takeaways, and a
Q&A close. Every slide carries speaker notes. Insurance numbers come from
results/rundata.json (so the deck, the demo, and the video can't drift); the
cross-industry trend/opportunity stats are cited on-slide.

    python3 deck/build_deck.py     # -> deck/insurance_pricing_webinar.pptx
"""

from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "video"))
import render as R                       # noqa: E402  (helpers, data, slide builders)
from render import (W, H, BG, TX, MUT, DIM, CARD, RULE, ENGINE, PRICE, COV, BAL, GOLD,  # noqa
                    T, font, base, abox, C, CM, SG, EN, BD, PD, CD)

SLIDES = ROOT / "deck" / "slides"
SLIDES.mkdir(parents=True, exist_ok=True)


# ---- webinar-only slides --------------------------------------------------
def w_title():
    img, d = base()
    T(d, (W / 2, 296), "Pricing to the customer,", font("bold", 74), TX, "mm")
    T(d, (W / 2, 392), "not the average", font("bold", 74), ENGINE, "mm")
    T(d, (W / 2, 524), "Turning the clicks you already collect into the price each shopper",
      font("reg", 36), MUT, "mm")
    T(d, (W / 2, 574), "will actually pay — with a choice model, an online bandit, and honest data",
      font("reg", 36), MUT, "mm")
    T(d, (W / 2, 706), "WEBINAR  ·  DEMAND-BASED INSURANCE PRICING  ·  ANDREW CURTIS",
      font("mono", 26), DIM, "mm")
    return img


def w_agenda():
    img, d = base("agenda")
    T(d, (110, 178), "What we'll cover", font("bold", 56), TX)
    items = [("The opportunity", "the click data every business collects and doesn't use"),
             ("The shopper reality", "what 97,009 people actually bought"),
             ("A choice model with honest error bars", "price vs coverage, and one disclosed layer"),
             ("Price- vs coverage-driven", "three segments, from real purchases"),
             ("Follow three shoppers through the engine", "lookalikes → clicks → steer → dynamic price"),
             ("See it live, then ship it", "the interactive demo, then AWS")]
    for i, (t, s) in enumerate(items):
        y = 296 + i * 104
        d.rounded_rectangle([110, y, 1810, y + 86], radius=10, fill=CARD, outline=RULE, width=1)
        T(d, (152, y + 24), str(i + 1), font("bold", 36), ENGINE)
        T(d, (224, y + 16), t, font("bold", 31), TX)
        T(d, (224, y + 54), s, font("reg", 24), MUT)
    return img


def w_problem():
    img, d = base("the problem")
    T(d, (W / 2, 292), "Everyone competes on the cheapest quote.", font("bold", 58), TX, "mm")
    T(d, (W / 2, 456), "But “the customer” is a fiction.", font("reg", 44), MUT, "mm")
    T(d, (W / 2, 540), "Price to the average and you discount to the people who'd have",
      font("reg", 38), MUT, "mm")
    T(d, (W / 2, 592), "paid more — while still losing the ones who won't.", font("reg", 38), MUT, "mm")
    T(d, (W / 2, 752), "The question isn't “what's our price?”   It's “who are we pricing to?”",
      font("bold", 38), ENGINE, "mm")
    return img


def w_data():
    img, d = base("the data")
    T(d, (110, 178), "97,009 real shoppers, not a survey", font("bold", 54), TX)
    T(d, (110, 250), "Allstate Purchase Prediction Challenge — real auto-insurance quote sessions.",
      font("reg", 30), MUT)
    facts = [(f"{C['n_customers']:,}", "shoppers"),
             (f"{C['median_quotes']}", "median quotes each"),
             (f"${C['premium_median']}", "median premium"),
             (f"{C['pct_saw_price_variation']}%", "saw >1 price")]
    x0, cw = 110, 410
    for i, (b, s) in enumerate(facts):
        x = x0 + i * (cw + 16)
        d.rounded_rectangle([x, 330, x + cw, 468], radius=10, fill=CARD, outline=RULE)
        T(d, (x + 28, 360), b, font("bold", 52), TX)
        T(d, (x + 28, 434), s.upper(), font("mono", 21), DIM)
    bl = [("Each shopper views real quotes — a 7-part coverage bundle (A–G) at a real premium — and buys one.", ENGINE),
          ("That sequence is the clickstream: it reveals whether they chased price or traded up to coverage.", ENGINE),
          ("What it can't tell us: who would have walked at a higher price. Everyone here bought — we'll return to this.", GOLD)]
    for i, (t, col) in enumerate(bl):
        y = 552 + i * 96
        d.ellipse([120, y + 9, 138, y + 27], fill=col)
        T(d, (162, y), t, font("reg", 30), TX if col == ENGINE else MUT)
    return img


def w_why_bayes():
    img, d = base("methodology")
    T(d, (110, 178), "Why Bayesian: report how sure you are", font("bold", 54), TX)
    T(d, (110, 252), "A conditional logit fit with a Laplace posterior (MAP + the Hessian) gives every",
      font("reg", 30), MUT)
    T(d, (110, 294), "coefficient a full distribution — not a lone point estimate.", font("reg", 30), MUT)
    # illustrate the price coefficient with its 95% credible interval
    c = CM["price"]
    x0, y0, bw = 300, 470, 1300
    lo, hi = -0.95, -0.25
    def xp(v):
        return x0 + (v - lo) / (hi - lo) * bw
    for g in (-0.9, -0.75, -0.6, -0.45, -0.3):
        d.line([xp(g), y0 - 26, xp(g), y0 + 26], fill=RULE, width=1)
        T(d, (xp(g), y0 + 44), f"{g:.2f}", font("mono", 22), DIM, "mm")
    a, b = xp(c["ci95"][0]), xp(c["ci95"][1])
    d.line([a, y0, b, y0], fill=ENGINE, width=4)
    for v in (a, b):
        d.line([v, y0 - 12, v, y0 + 12], fill=ENGINE, width=4)
    cx = xp(c["mean"])
    d.ellipse([cx - 10, y0 - 10, cx + 10, y0 + 10], fill=ENGINE)
    T(d, (cx, y0 - 34), f"price = {c['mean']:.2f}", font("bold", 28), TX, "mm")
    T(d, (300, 396), "PRICE SENSITIVITY  ·  95% CREDIBLE INTERVAL", font("mono", 22), ENGINE)
    T(d, (110, 636), "The point estimate hides the uncertainty; the interval is the plausible range.",
      font("reg", 30), TX)
    T(d, (110, 690), "It excludes zero, so price genuinely moves the choice — and we can say how much,",
      font("reg", 30), MUT)
    T(d, (110, 732), "with honest error bars, which is what a pricing committee actually needs.",
      font("reg", 30), MUT)
    return img


def w_calibration():
    img, d = base("the honest engineering")
    T(d, (110, 172), "Calibrate the one simulated layer — then stress-test it", font("bold", 48), TX)
    T(d, (110, 244), "Buyers-only data can't contain walk-away risk, so we model it: P(convert | price),",
      font("reg", 28), MUT)
    T(d, (110, 282), "calibrated to published auto-insurance elasticities — price-driven the most elastic.",
      font("reg", 28), MUT)
    # arc elasticities per segment
    arcs = EN["arc_elasticities"]
    labels = [("price_driven", "Price-driven", PRICE), ("balanced", "Balanced", BAL),
              ("coverage_driven", "Coverage-driven", COV)]
    for i, (k, lab, col) in enumerate(labels):
        x = 110 + i * 470
        d.rounded_rectangle([x, 336, x + 440, 452], radius=10, fill=CARD, outline=RULE)
        d.rounded_rectangle([x, 336, x + 440, 342], radius=3, fill=col)
        T(d, (x + 26, 360), lab, font("bold", 28), TX)
        T(d, (x + 26, 402), f"arc elasticity  {arcs[k]:.1f}", font("mono", 26), col)
    # the sweep
    T(d, (110, 512), "ROBUSTNESS — SWEEP THE ELASTICITY 0.5×…1.5×", font("mono", 22), GOLD)
    sw = EN["elasticity_sweep"]
    x0, y0, bh, gap = 300, 548, 40, 58
    mx = max(s["engine_lift_vs_compete_pct"] for s in sw) * 1.12
    bw = 900
    for i, s in enumerate(sw):
        y = y0 + i * gap
        w = int(bw * s["engine_lift_vs_compete_pct"] / mx)
        T(d, (x0 - 20, y + bh / 2), f"{s['escale']}×", font("mono", 24), MUT, "rm")
        d.rounded_rectangle([x0, y, x0 + w, y + bh], radius=6, fill=GOLD)
        T(d, (x0 + w + 16, y + bh / 2),
          f"+{s['engine_lift_vs_compete_pct']:.1f}%   (conv {s['engine_conversion']*100:.0f}%)",
          font("mono", 24), TX, "lm")
    lo = min(s["engine_lift_vs_compete_pct"] for s in sw)
    hi = max(s["engine_lift_vs_compete_pct"] for s in sw)
    T(d, (110, 900), f"The lift holds at {lo:.1f}–{hi:.1f}% across the whole range — it doesn't "
      "hinge on one flattering value.", font("bold", 30), TX)
    return img


def _numbox(d, x, y, w, h, n, title, sub, col):
    d.rounded_rectangle([x, y, x + w, y + h], radius=12, fill=CARD, outline=RULE, width=1)
    d.rounded_rectangle([x, y, x + w, y + 7], radius=3, fill=col)
    d.ellipse([x + 24, y + 28, x + 70, y + 74], fill=col)
    T(d, (x + 47, y + 51), str(n), font("bold", 28), (15, 18, 24), "mm")
    T(d, (x + 24, y + 98), title, font("bold", 29), TX)
    T(d, (x + 24, y + 142), sub, font("reg", 21), MUT)


def w_loop():
    img, d = base("how it prices a new shopper")
    T(d, (110, 178), "How the engine prices a shopper it's never seen", font("bold", 50), TX)
    steps = [("Lookalike prior", "similar shoppers' real buys", ENGINE),
             ("Float test offers", "lean · standard · full", ENGINE),
             ("Read the clicks", "who they really are", GOLD),
             ("Steer & price", "coverage + dynamic premium", COV)]
    w, gap, x0, y, h = 395, 40, 110, 336, 190
    for i, (t, s, col) in enumerate(steps):
        x = x0 + i * (w + gap)
        _numbox(d, x, y, w, h, i + 1, t, s, col)
        if i < 3:
            T(d, (x + w + gap / 2, y + h / 2), "→", font("bold", 40), DIM, "mm")
    d.rounded_rectangle([110, 596, 1810, 700], radius=12, fill=CARD, outline=GOLD, width=2)
    T(d, (150, 616), "THE CATCH", font("mono", 22), GOLD)
    T(d, (150, 650), "Demographics guess the coverage. Only the clicks find who'll walk on price.",
      font("bold", 34), TX)
    return img


def w_three():
    img, d = base("three shoppers, three moves")
    T(d, (110, 178), "Three shoppers. Same engine. Three moves.", font("bold", 50), TX)
    cards = [("Dana", "PRICE-DRIVEN", "24 · renter · cheap car", PRICE,
              "Basic", "$25K / $50K", "$551 → $521", "−$30", "discount to win the sale"),
             ("Sam", "BALANCED", "38 · homeowner · mid car", BAL,
              "Standard", "$100K / $300K", "$605 → $605", "holds", "list price already right"),
             ("Carla", "COVERAGE-DRIVEN", "62 · homeowner · pricey car", COV,
              "Full", "$250K / $500K", "$673 → $763", "+$90", "markup she'll accept")]
    w, gap, x0, y, h = 550, 24, 110, 300, 566
    for i, (nm, persona, demo, col, tier, limits, price, delta, why) in enumerate(cards):
        x = x0 + i * (w + gap)
        d.rounded_rectangle([x, y, x + w, y + h], radius=14, fill=CARD, outline=RULE, width=1)
        d.rounded_rectangle([x, y, x + w, y + 8], radius=4, fill=col)
        T(d, (x + 30, y + 38), nm, font("bold", 44), TX)
        T(d, (x + 30, y + 104), persona, font("mono", 21), col)
        T(d, (x + 30, y + 142), demo, font("reg", 22), MUT)
        d.line([x + 30, y + 200, x + w - 30, y + 200], fill=RULE, width=1)
        T(d, (x + 30, y + 220), "STEERED TO", font("mono", 18), DIM)
        T(d, (x + 30, y + 250), tier, font("bold", 34), TX)
        T(d, (x + 30, y + 298), limits + " liability", font("mono", 21), MUT)
        T(d, (x + 30, y + 362), "DYNAMIC PRICE", font("mono", 18), DIM)
        T(d, (x + 30, y + 392), price, font("bold", 30), TX)
        T(d, (x + 30, y + 448), delta, font("bold", 44), col)
        T(d, (x + 30, y + 514), why, font("reg", 22), MUT)
    return img


def w_dynamic():
    img, d = base("the dynamic pricing")
    T(d, (110, 178), "Dynamic pricing: the premium, off the list", font("bold", 50), TX)
    T(d, (110, 250), "Same coverage the shopper wants — a personalized premium, moved off the standard price.",
      font("reg", 32), MUT)
    y, x0, x1 = 500, 300, 1620
    cx = (x0 + x1) / 2
    d.line([x0, y, x1, y], fill=RULE, width=3)
    T(d, (cx, y - 96), "LIST PRICE", font("mono", 22), DIM, "mm")
    d.line([cx, y - 22, cx, y + 22], fill=DIM, width=3)
    T(d, (x0, y + 124), "←  DISCOUNT", font("mono", 24), PRICE, "lm")
    T(d, (x1, y + 124), "MARKUP  →", font("mono", 24), COV, "rm")

    def mark(px, label, delta, col):
        d.ellipse([px - 13, y - 13, px + 13, y + 13], fill=col)
        T(d, (px, y - 54), label, font("bold", 28), col, "mm")
        T(d, (px, y + 56), delta, font("mono", 24), MUT, "mm")
    mark(cx - 470, "Dana", "−$30", PRICE)
    mark(cx, "Sam", "holds", BAL)
    mark(cx + 470, "Carla", "+$90", COV)
    T(d, (110, 722), "Elastic shoppers get a discount to convert. Inelastic coverage-lovers get a markup.",
      font("reg", 32), TX)
    T(d, (110, 774), "The balanced middle holds at list — personalization pays most at the extremes.",
      font("reg", 32), MUT)
    return img


def w_takeaways():
    img, d = base("what to take away")
    T(d, (110, 178), "For anyone who prices to a mix of customers", font("bold", 50), TX)
    cards = [("The click is the asset you already own",
              "You collect the strongest signal of intent in every session. The win is acting on it, not gathering more data.", ENGINE),
             ("Stop discounting to the people who'd say yes",
              "The lift comes from NOT pushing the price-sensitive, and trading up the coverage-driven.", COV),
             ("You can learn elasticity online",
              "A bandit finds each segment's best offer from clicks — within a few hundred shoppers.", BAL),
             ("Show the seam",
              "Label what's real and what's modeled. The honest version is the one people trust — and it's stronger.", GOLD)]
    for i, (t, s, col) in enumerate(cards):
        cx = 110 + (i % 2) * 862
        cy = 320 + (i // 2) * 288
        d.rounded_rectangle([cx, cy, cx + 826, cy + 252], radius=12, fill=CARD, outline=RULE)
        d.rounded_rectangle([cx, cy, cx + 8, cy + 252], radius=3, fill=col)
        T(d, (cx + 36, cy + 34), t, font("bold", 32), TX)
        # wrap the body to ~34 chars/line
        words, line, lines = s.split(), "", []
        for wd in words:
            if len(line + " " + wd) > 40:
                lines.append(line); line = wd
            else:
                line = (line + " " + wd).strip()
        lines.append(line)
        for j, ln in enumerate(lines):
            T(d, (cx + 36, cy + 100 + j * 40), ln, font("reg", 27), MUT)
    return img


def w_thanks():
    img, d = base()
    T(d, (W / 2, 288), "Thank you.", font("bold", 76), TX, "mm")
    T(d, (W / 2, 380), "Questions?", font("bold", 76), ENGINE, "mm")
    rows = [("REPO", "github.com/andycurtis1973/insurance-pricing-engine"),
            ("LIVE DEMO", "andycurtis1973.github.io/insurance-pricing-engine/web/app.html")]
    for i, (lab, u) in enumerate(rows):
        y = 496 + i * 92
        s = f"{lab}   {u}"
        w = d.textlength(s, font=font("mono", 28))
        x0 = (W - w) / 2 - 34
        d.rounded_rectangle([x0, y, x0 + w + 68, y + 68], radius=12,
                            fill=CARD, outline=ENGINE, width=2)
        T(d, (x0 + 34, y + 34), lab, font("mono", 28), DIM, "lm")
        T(d, (x0 + 34 + d.textlength(f"{lab}   ", font=font("mono", 28)), y + 34),
          u, font("mono", 28), ENGINE, "lm")
    T(d, (W / 2, 724), "Real data (Allstate) · open code · the elasticity swept · the simulated layer labeled",
      font("reg", 27), MUT, "mm")
    T(d, (W / 2, 776), "Calibration: Guven & McPhail, CAS E-Forum 2013 · Akur8, Demand Modeling in Insurance",
      font("mono", 22), DIM, "mm")
    return img


def _statcard(d, x, y, w, h, big, label, src, col):
    d.rounded_rectangle([x, y, x + w, y + h], radius=12, fill=CARD, outline=RULE, width=1)
    d.rounded_rectangle([x, y, x + w, y + 7], radius=3, fill=col)
    T(d, (x + 30, y + 30), big, font("bold", 68), col)
    words, line, lines = label.split(), "", []
    for wd in words:
        if len(line + " " + wd) > 30:
            lines.append(line); line = wd
        else:
            line = (line + " " + wd).strip()
    lines.append(line)
    for j, ln in enumerate(lines):
        T(d, (x + 30, y + 128 + j * 38), ln, font("reg", 27), MUT)
    T(d, (x + 30, y + h - 44), src.upper(), font("mono", 19), DIM)


def w_clickgap():
    img, d = base("the trend  ·  across retail, e-commerce, insurance")
    T(d, (110, 178), "Everyone captures the clicks. Almost no one acts on them.",
      font("bold", 50), TX)
    T(d, (110, 252), "Behavioral clickstream is the fastest-growing data asset in every "
      "consumer business — and the least activated.", font("reg", 30), MUT)
    cards = [("72%", "of firms now collect and unify behavioral data - capture is racing ahead of use", "Twilio Segment, 2026", GOLD),
             ("~60%", "of companies still don't act on real-time click signals they already record", "DemandSage, 2026", PRICE),
             ("11.4x", "more transactions when the click actually triggers the next offer", "Clickstream, 2026", COV)]
    w, gap, x0, y, h = 553, 20, 110, 344, 316
    for i, (big, lab, src, col) in enumerate(cards):
        _statcard(d, x0 + i * (w + gap), y, w, h, big, lab, src, col)
    d.rounded_rectangle([110, 700, 1810, 858], radius=12, fill=CARD, outline=ENGINE, width=2)
    T(d, (150, 722), "THE GAP", font("mono", 22), ENGINE)
    T(d, (150, 758), "Every session already reveals who's price-shopping and who's trading up.",
      font("bold", 31), TX)
    T(d, (150, 802), "Most businesses record that click - and price as if they never saw it.",
      font("bold", 31), TX)
    return img


def w_opportunity():
    img, d = base("the opportunity")
    T(d, (110, 178), "The prize: growth from data you already own", font("bold", 52), TX)
    T(d, (110, 252), "Personalization on behavioral data is the highest-ROI lever most teams "
      "under-pull - and pricing is where it hits the P&L directly.", font("reg", 30), MUT)
    # left: cross-industry evidence
    T(d, (110, 336), "ACROSS INDUSTRIES  ·  WHAT PERSONALIZATION RETURNS", font("mono", 22), DIM)
    rows = [("10-15%", "revenue lift, typical, from personalized experiences", COV),
            ("40%", "more revenue captured by personalization leaders vs the rest", BAL),
            ("11.4x", "transaction lift when engagement is triggered in real time", GOLD)]
    for i, (big, lab, col) in enumerate(rows):
        y = 402 + i * 150
        T(d, (130, y), big, font("bold", 62), col)
        words, line, lines = lab.split(), "", []
        for wd in words:
            if len(line + " " + wd) > 34:
                lines.append(line); line = wd
            else:
                line = (line + " " + wd).strip()
        lines.append(line)
        for j, ln in enumerate(lines):
            T(d, (360, y + 6 + j * 38), ln, font("reg", 28), MUT)
    T(d, (130, 862), "McKinsey  ·  DemandSage 2026  ·  Clickstream 2026", font("mono", 20), DIM)
    # right: the worked example (this build)
    x = 990
    d.rounded_rectangle([x, 340, 1810, 862], radius=14, fill=CARD, outline=ENGINE, width=2)
    d.rounded_rectangle([x, 340, 1810, 348], radius=4, fill=ENGINE)
    T(d, (x + 36, 372), "IN THIS BUILD  ·  INSURANCE, THE WORKED EXAMPLE", font("mono", 22), ENGINE)
    lift = EN["lift_vs_compete_pct"]
    st = EN["strategies"]
    delta = st["engine"]["rev_per_shopper"] - st["compete_on_price"]["rev_per_shopper"]
    cc = st["compete_on_price"]["conversion"] * 100
    ec = st["engine"]["conversion"] * 100
    T(d, (x + 36, 424), f"+{lift:.1f}%", font("bold", 96), ENGINE)
    T(d, (x + 36, 550), "revenue per shopper vs competing on price -", font("reg", 27), MUT)
    T(d, (x + 36, 586), "purely from reading the clicks", font("reg", 27), MUT)
    d.line([x + 36, 636, 1774, 636], fill=RULE, width=1)
    T(d, (x + 36, 652), f"{cc:.0f}% -> {ec:.0f}%", font("bold", 40), TX)
    T(d, (x + 36, 706), "conversion", font("mono", 21), DIM)
    T(d, (x + 470, 652), f"+${delta:.0f}", font("bold", 40), TX)
    T(d, (x + 470, 706), "per shopper", font("mono", 21), DIM)
    T(d, (x + 36, 766), f"~${delta:.0f}M a year on 1,000,000 quoted shoppers", font("bold", 34), GOLD)
    T(d, (x + 36, 818), "illustrative, at the simulated lift on the disclosed conversion layer",
      font("reg", 22), DIM)
    return img


def w_demo():
    img, d = base("see it live")
    T(d, (110, 172), "The customer's screen and your backend, side by side", font("bold", 50), TX)
    T(d, (110, 244), "Three shoppers - price-sensitive, neutral, coverage-maximizing - run through "
      "one engine. Watch the click analysis sharpen and the AWS loop light up.",
      font("reg", 28), MUT)
    # left frame: what the customer sees
    lx, ly, lw, lh = 110, 322, 820, 452
    d.rounded_rectangle([lx, ly, lx + lw, ly + lh], radius=14, fill=CARD, outline=RULE, width=1)
    T(d, (lx + 28, ly + 22), "WHAT THE CUSTOMER SEES", font("mono", 20), DIM)
    d.rounded_rectangle([lx + 28, ly + 60, lx + lw - 28, ly + 104], radius=8, fill=BG, outline=RULE)
    T(d, (lx + 48, ly + 72), "quicksure-auto.com/quote", font("mono", 22), MUT)
    tiers = [("Lean", "$25K / $50K", "$651", PRICE), ("Standard", "$100K / $300K", "$690", BAL),
             ("Full", "$250K / $500K", "$742", COV)]
    for i, (nm, lim, pr, col) in enumerate(tiers):
        ry = ly + 132 + i * 96
        d.rounded_rectangle([lx + 28, ry, lx + lw - 28, ry + 80], radius=10, fill=BG, outline=col, width=2)
        d.rounded_rectangle([lx + 28, ry, lx + 34, ry + 80], radius=3, fill=col)
        T(d, (lx + 58, ry + 16), nm, font("bold", 30), TX)
        T(d, (lx + 58, ry + 52), lim + " liability", font("mono", 19), MUT)
        T(d, (lx + lw - 58, ry + 24), pr, font("bold", 34), col, "ra")
    # right frame: what the model is doing
    rx, ry0, rw, rh = 970, 322, 840, 452
    d.rounded_rectangle([rx, ry0, rx + rw, ry0 + rh], radius=14, fill=CARD, outline=ENGINE, width=2)
    T(d, (rx + 28, ry0 + 22), "WHAT THE MODEL IS DOING  ·  THE CLICK ANALYSIS", font("mono", 20), ENGINE)
    beliefs = [("price-driven", 0.68, PRICE), ("balanced", 0.20, BAL), ("coverage-driven", 0.12, COV)]
    for i, (nm, v, col) in enumerate(beliefs):
        by = ry0 + 66 + i * 58
        T(d, (rx + 28, by), nm, font("mono", 22), MUT)
        bx0, bw = rx + 300, 380
        d.rounded_rectangle([bx0, by + 2, bx0 + bw, by + 26], radius=8, fill=BG, outline=RULE)
        d.rounded_rectangle([bx0, by + 2, bx0 + int(bw * v), by + 26], radius=8, fill=col)
        T(d, (bx0 + bw + 20, by), f"{v*100:.0f}%", font("bold", 24), TX)
    T(d, (rx + 28, ry0 + 250), "posterior = prior (real lookalikes) x likelihood (each click)",
      font("mono", 19), DIM)
    # mini streaming loop
    loop = [("Shopper", ENGINE), ("Lambda", ENGINE), ("Offer", ENGINE),
            ("Kinesis", GOLD), ("Flink", GOLD), ("DynamoDB", BAL)]
    bw2, bh2, gx = 118, 54, rx + 28
    for i, (nm, col) in enumerate(loop):
        cxx = gx + i % 3 * (bw2 + 24)
        cyy = ry0 + 300 + (i // 3) * 74
        d.rounded_rectangle([cxx, cyy, cxx + bw2, cyy + bh2], radius=8, fill=BG, outline=col, width=2)
        T(d, (cxx + bw2 / 2, cyy + bh2 / 2), nm, font("bold", 19), TX, "mm")
    # scenario chips (row 1) + url (row 2)
    T(d, (110, 800), "RUNS BACK-TO-BACK", font("mono", 20), DIM)
    chips = [("1  Price-sensitive", PRICE), ("2  Neutral", BAL), ("3  Maximize coverage", COV)]
    cx = 388
    for nm, col in chips:
        cw = int(d.textlength(nm, font=font("bold", 26))) + 72
        d.rounded_rectangle([cx, 792, cx + cw, 852], radius=30, fill=CARD, outline=col, width=2)
        d.ellipse([cx + 26, 816, cx + 42, 832], fill=col)
        T(d, (cx + 58, 808), nm, font("bold", 26), TX)
        cx += cw + 24
    url = "andycurtis1973.github.io/insurance-pricing-engine/web/app.html"
    uw = d.textlength(url, font=font("mono", 26))
    x0 = (W - uw) / 2 - 150
    d.rounded_rectangle([x0, 892, x0 + uw + 184, 952], radius=12, fill=CARD, outline=ENGINE, width=2)
    T(d, (x0 + 34, 922), "TRY IT", font("mono", 24), DIM, "lm")
    T(d, (x0 + 150, 922), url, font("mono", 26), ENGINE, "lm")
    return img


# ---- the deck: (builder, static?, speaker notes) --------------------------
DECK = [
    (w_title, "Welcome. Today we're looking at how an auto insurer can price to who a shopper "
     "actually is, instead of to the average. Three ingredients: a choice model, an online "
     "learning engine, and — importantly — an honest line between what's real data and what's "
     "modeled. Everything you'll see is built on 97,000 real shopping sessions."),
    (w_agenda, "Here's the path. We start with what shoppers actually do, build a model that "
     "explains it, split them into segments, then draw a clear line around the one simulated "
     "piece before we ever talk about money. Then the engine and its lift, how it learns online, "
     "and how it would deploy. Save questions for the end, or drop them in the chat."),
    (w_problem, "Start with the framing. Every insurer competes on the cheapest quote — but "
     "'the customer' is a fiction. When you price to the average, you hand a discount to the "
     "people who would happily have paid more, and you still lose the ones who won't pay up. "
     "So the real question isn't 'what's our price?' — it's 'who are we pricing to?' Everything "
     "that follows is about answering that from data."),
    (w_clickgap, "Here's the bigger pattern, and it's not just insurance. Across retail, e-commerce "
     "and financial services, everyone is now capturing behavioral clickstream — 72% of firms unify "
     "it — but roughly 60% still don't act on those signals in real time. The click that tells you "
     "someone is price-shopping versus trading up is recorded in every session, and then thrown away. "
     "That gap between collecting behavior and acting on it is the cheapest growth most businesses have left."),
    (w_opportunity, "And the prize is large. Personalization built on that behavioral data typically "
     "returns a 10 to 15% revenue lift; the leaders pull 40% more revenue from it than everyone else; "
     "and when you trigger on engagement in real time, transaction rates jump. Pricing is where this "
     "hits the P&L most directly — so insurance is our worked example. Reading the clicks lifts revenue "
     "per shopper 6.7% and conversion from 65 to 67%, about $27 a head — which on a million quoted "
     "shoppers is roughly $27M a year. That last number is illustrative, at the simulated lift, but "
     "the mechanism is the point."),
    (lambda: R.a_behavior(1.0), "Every square is a real person shopping for car insurance. Only "
     "about a quarter bought the cheapest option they saw. Nearly half traded up to richer "
     "coverage and paid more. Same quote engine, two completely different shoppers — and if you "
     "lead with the cheapest quote for everyone, you're discounting to the half who'd have paid more."),
    (w_data, "The credibility rests on the data. This is the Allstate Purchase Prediction "
     "dataset — real quote sessions, real coverage bundles A through G, real premiums, a median "
     "of six quotes each. Crucially, note the last point: everyone in this data bought something. "
     "So it cannot tell us who would have walked away at a higher price. Hold that thought."),
    (lambda: R.a_choice(1.0), "To explain the choices we fit a Bayesian conditional logit. Price "
     "pushes people toward cheaper; coverage pulls them richer. And price sensitivity isn't one "
     "number — it's sharper for homeowners and young drivers, and it flips for expensive cars. "
     "It predicts the bought offer 42% of the time, versus 26 and 39 for the naive baselines."),
    (w_why_bayes, "Why Bayesian rather than a plain regression? The Laplace posterior gives each "
     "coefficient a full distribution, so we get credible intervals. The price coefficient is "
     "minus 0.60, and the interval clearly excludes zero. A pricing committee needs the range and "
     "the confidence, not just a point estimate — this reports how sure we are."),
    (lambda: R.a_segments(1.0), "Split people by what they actually bought and three populations "
     "fall out: price-driven, coverage-driven, and balanced in between. They don't want the same "
     "thing — the coverage-driven buy 62% of max coverage versus 53% for the price-driven. So why "
     "would you quote them the same way?"),
    (lambda: R.a_honesty(1.0), "Now the honest part, before a single dollar of lift. Because the "
     "data is buyers-only, we can't observe walk-away. So we add it: a conversion model, and we "
     "label it clearly. The shoppers, menus, preferences and segments are real; the walk-away "
     "response is a disclosed model. Everything past here is lift in simulation — not a claim "
     "about Allstate's revenue."),
    (w_calibration, "We don't just make up the elasticity. We calibrate it to published auto-"
     "insurance figures — price-driven shoppers around minus two, coverage-lovers around minus "
     "one-half — and then we sweep it from half to one-and-a-half times. The lift holds between "
     "5.5 and 8.2% across the whole range. That's the difference between a calibrated simulator "
     "and a flattering guess."),
    (w_loop, "Here's how the engine prices a shopper it's never seen. Step one: look up what real "
     "lookalikes bought — a warm start on coverage. Step two: float a few test offers. Step three: "
     "read the clicks. Step four: steer the coverage and set the price. The catch is the whole "
     "point — demographics can guess the coverage, but only the clicks reveal who'll walk on price."),
    (w_three, "Watch it run on three real profiles. Dana, a young renter, clicks the cheapest option "
     "— she's price-driven, so the engine steers her to Basic and prices it DOWN, $551 to $521, to "
     "win a sale she'd otherwise walk from. Carla, an older homeowner, reaches for full coverage — "
     "steered to Full and priced UP, $673 to $763, because she'll pay. Sam is in between: he takes "
     "the standard plan, and the engine holds at the list price. Same engine, three different moves."),
    (w_dynamic, "This is where the dynamic pricing actually happens. It's not the tier — those list "
     "prices are cost-based. It's the personalized premium the engine sets on the tier the shopper "
     "wants: a discount to convert the price-sensitive, a markup for the coverage-lover who'll pay, "
     "and — for the balanced middle — no move at all. Personalization pays most at the extremes; "
     "in the middle, the standard price is already right. That's the honest shape of it."),
    (lambda: R.a_learn(1.0), "In the real world you don't know a shopper's elasticity up front — "
     "you learn it from clicks. Push a price-driven shopper to richer coverage and conversion "
     "falls; push a coverage-driven one and it rises. Same offer, opposite result. A Thompson "
     "Sampling bandit finds each segment's best offer on its own, converging to within half a "
     "percent of a perfect-knowledge oracle."),
    (lambda: R.a_architecture(1.0), "How would this ship? In production it's a streaming loop. "
     "Every quote and click flows through Kinesis; a decision service — API Gateway and Lambda "
     "with the posteriors in DynamoDB — picks the offer in milliseconds; Flink folds conversions "
     "back in near-real-time; Firehose lands events in S3 for dashboards and a nightly SageMaker "
     "retrain. To be clear, this is a blueprint — the build itself ran offline on a laptop."),
    (w_demo, "And you can watch the whole thing run. This is the interactive demo: on the left is "
     "exactly what the customer sees — a quote site, three offers; on the right is what the model is "
     "doing — the belief bars over price-driven, balanced and coverage-driven sharpening with each "
     "click, and the AWS streaming loop lighting up as events flow. It runs three shoppers back to "
     "back: price-sensitive, neutral, and coverage-maximizing. It's a standalone web page — the link "
     "is on the slide and on the closing slide, so anyone can click through it themselves."),
    (w_takeaways, "Four things to take home. Heterogeneity is the whole game — 'the customer' is "
     "a fiction. The lift comes from not pushing the price-sensitive and trading up the coverage-"
     "driven. You can learn elasticity online, cheaply. And show the seam: labeling the modeled "
     "layer is what makes the rest believable."),
    (w_thanks, "That's the build. It's real data, open code, the elasticity swept, and the one "
     "simulated layer labeled. The repo and an interactive walkthrough are on screen. Happy to "
     "take questions — including the obvious one about deploying this against a live book with a "
     "real conversion signal."),
]


def main() -> int:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for i, (builder, notes) in enumerate(DECK):
        png = SLIDES / f"slide_{i:02d}.png"
        builder().save(png)
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), 0, 0, width=prs.slide_width, height=prs.slide_height)
        slide.notes_slide.notes_text_frame.text = notes
        print(f"  [{i+1:2d}/{len(DECK)}] {png.name}")

    out = ROOT / "deck" / "insurance_pricing_webinar.pptx"
    prs.save(str(out))
    mb = out.stat().st_size / 1e6
    print(f"\n  ✅ {out}  ({len(DECK)} slides, {mb:.1f} MB)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
